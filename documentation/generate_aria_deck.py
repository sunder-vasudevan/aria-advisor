#!/usr/bin/env python3
"""
Generates a McKinsey/PwC-style 12-slide executive deck for the ARIA Platform.
Run:    python generate_aria_deck.py
Output: documentation/ARIA_Executive_Deck.pptx
Requires: pip install python-pptx matplotlib
"""

import io
import os
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
DARK_NAVY  = RGBColor(0x00, 0x2B, 0x5C)
MID_BLUE   = RGBColor(0x00, 0x5B, 0x99)
BRAND_BLUE = RGBColor(0x1D, 0x6F, 0xDB)   # ARIA #1D6FDB
TEAL       = RGBColor(0x00, 0x7A, 0x87)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF4, 0xF7)
DARK_TEXT  = RGBColor(0x1A, 0x1A, 0x2E)
MID_GREY   = RGBColor(0x88, 0x88, 0x88)
AMBER      = RGBColor(0xE8, 0x8C, 0x00)
ROSE       = RGBColor(0xE5, 0x3E, 0x3E)
SLATE_50   = RGBColor(0xF8, 0xFA, 0xFC)

# ── Slide dimensions (widescreen 16:9) ────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ── Low-level helpers ─────────────────────────────────────────────────────────

def blank_slide(prs):
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


def rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def txbox(slide, left, top, width, height,
          text='', size=18, bold=False, italic=False,
          color=DARK_TEXT, align=PP_ALIGN.LEFT,
          wrap=True, line_spacing=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        from pptx.oxml.ns import qn as _qn
        from lxml import etree
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, _qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, _qn('a:spcPct'))
        spcPct.set('val', str(int(line_spacing * 1000)))
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    return box, tf


def footer_bar(slide, label='A-RiA · ARIA Personal  |  March 2026  |  Made with \u2764\ufe0f in Hyderabad'):
    rect(slide, 0, H - Inches(0.3), W, Inches(0.3), fill_color=DARK_NAVY)
    txbox(slide, Inches(0.3), H - Inches(0.28), W - Inches(0.6), Inches(0.26),
          text=label, size=Pt(7), color=WHITE, align=PP_ALIGN.LEFT)


def section_label(slide, text, left=Inches(0.5), top=Inches(0.18)):
    rect(slide, left, top, Inches(2.8), Inches(0.22), fill_color=TEAL)
    txbox(slide, left + Inches(0.08), top, Inches(2.6), Inches(0.22),
          text=text.upper(), size=Pt(7.5), bold=True, color=WHITE)


def slide_title(slide, text, top=Inches(0.52), color=DARK_NAVY):
    txbox(slide, Inches(0.5), top, Inches(12.3), Inches(0.6),
          text=text, size=Pt(26), bold=True, color=color)


def rule_line(slide, top, color=MID_BLUE, left=Inches(0.5), width=None):
    w = width or (W - Inches(1.0))
    rect(slide, left, top, w, Inches(0.025), fill_color=color)


def big_number(slide, number, label, left, top, num_color=BRAND_BLUE):
    txbox(slide, left, top, Inches(2.8), Inches(1.1),
          text=number, size=Pt(48), bold=True, color=num_color, align=PP_ALIGN.CENTER)
    txbox(slide, left, top + Inches(1.0), Inches(2.8), Inches(0.35),
          text=label, size=Pt(10), color=MID_GREY, align=PP_ALIGN.CENTER)


def bullet_block(slide, items, left, top, width, size=Pt(11), color=DARK_TEXT,
                 bullet_char='\u25b8', spacing=Inches(0.32)):
    for i, item in enumerate(items):
        txbox(slide, left, top + i * spacing, width, Inches(0.3),
              text=f'{bullet_char}  {item}', size=size, color=color)


def callout(slide, text, left, top, width, height,
            bg=RGBColor(0xE8, 0xF4, 0xF6), border=TEAL):
    rect(slide, left, top, width, height, fill_color=bg, line_color=border, line_width=Pt(2))
    txbox(slide, left + Inches(0.15), top + Inches(0.12),
          width - Inches(0.3), height - Inches(0.24),
          text=text, size=Pt(11), italic=True, color=DARK_NAVY, wrap=True)


def two_col_table(slide, headers, rows, left, top, col_widths, row_height=Inches(0.32)):
    n_cols = len(headers)
    x = left
    for i, h in enumerate(headers):
        rect(slide, x, top, col_widths[i], row_height, fill_color=MID_BLUE)
        txbox(slide, x + Inches(0.08), top + Inches(0.05),
              col_widths[i] - Inches(0.1), row_height - Inches(0.08),
              text=h, size=Pt(9), bold=True, color=WHITE)
        x += col_widths[i]
    for r_idx, row in enumerate(rows):
        bg = LIGHT_GREY if r_idx % 2 == 0 else WHITE
        x = left
        for c_idx, cell in enumerate(row):
            rect(slide, x, top + (r_idx + 1) * row_height,
                 col_widths[c_idx], row_height, fill_color=bg)
            txbox(slide, x + Inches(0.08),
                  top + (r_idx + 1) * row_height + Inches(0.05),
                  col_widths[c_idx] - Inches(0.1), row_height - Inches(0.08),
                  text=str(cell), size=Pt(9), color=DARK_TEXT)
            x += col_widths[c_idx]


# ── Chart helpers ─────────────────────────────────────────────────────────────

def _c(rgb_tuple):
    if isinstance(rgb_tuple, RGBColor):
        return (rgb_tuple.red / 255, rgb_tuple.green / 255, rgb_tuple.blue / 255)
    return tuple(c / 255 for c in rgb_tuple)


_NAVY  = (0x00, 0x2B, 0x5C)
_BLUE  = (0x1D, 0x6F, 0xDB)
_TEAL  = (0x00, 0x7A, 0x87)
_AMBER = (0xE8, 0x8C, 0x00)
_ROSE  = (0xE5, 0x3E, 0x3E)
_LGREY = (0xF2, 0xF4, 0xF7)


def _savefig(fig, dpi=150):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight', facecolor='none')
    buf.seek(0)
    plt.close(fig)
    return buf


def chart_probability_pills():
    """Bar chart showing probability pill thresholds."""
    fig, ax = plt.subplots(figsize=(5, 2.8), facecolor='none')
    fig.patch.set_alpha(0)
    ax.set_facecolor('none')
    cats = ['At Risk\n(<40%)', 'Needs Attention\n(40–70%)', 'On Track\n(\u226570%)']
    vals = [25, 55, 82]
    colours = [_c(_ROSE), _c(_AMBER), _c(_TEAL)]
    bars = ax.bar(cats, vals, color=colours, edgecolor='white', linewidth=1.0, width=0.5)
    for bar, val in zip(bars, vals):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 1.5,
                f'{val}%', ha='center', va='bottom', fontsize=9, fontweight='bold',
                color=_c(_NAVY))
    ax.set_ylim(0, 100)
    ax.set_title('Goal Probability — Pill System', fontsize=10, fontweight='bold', color=_c(_NAVY), pad=8)
    for sp in ['top', 'right']:
        ax.spines[sp].set_visible(False)
    ax.spines['left'].set_color('#AAAAAA')
    ax.spines['bottom'].set_color('#AAAAAA')
    ax.tick_params(labelsize=8)
    fig.tight_layout()
    return _savefig(fig)


def chart_monte_carlo():
    """Illustrative Monte Carlo paths."""
    np.random.seed(42)
    fig, ax = plt.subplots(figsize=(5.5, 3.0), facecolor='none')
    fig.patch.set_alpha(0)
    ax.set_facecolor('none')
    months = np.arange(0, 121)
    for i in range(30):
        monthly_r = np.random.normal(0.01, 0.004, 120)
        vals = [500000]
        for r in monthly_r:
            vals.append(vals[-1] * (1 + r) + 10000)
        colour = _c(_TEAL) if vals[-1] >= 2500000 else _c(_ROSE)
        ax.plot(months, vals, color=colour, alpha=0.3, linewidth=0.9)
    ax.axhline(y=2500000, color=_c(_NAVY), linewidth=1.5, linestyle='--')
    ax.text(122, 2500000, 'Target', fontsize=8, color=_c(_NAVY), va='center')
    ax.set_xlabel('Months', fontsize=8)
    ax.set_ylabel('Portfolio Value (\u20b9)', fontsize=8)
    ax.set_title('1,000 Monte Carlo Paths', fontsize=10, fontweight='bold', color=_c(_NAVY), pad=8)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f'\u20b9{x/100000:.0f}L'))
    ax.tick_params(labelsize=7)
    for sp in ['top', 'right']:
        ax.spines[sp].set_visible(False)
    ax.spines['left'].set_color('#AAAAAA')
    ax.spines['bottom'].set_color('#AAAAAA')
    teal_patch = mpatches.Patch(color=_c(_TEAL), label='Success path')
    rose_patch = mpatches.Patch(color=_c(_ROSE), label='Miss path')
    ax.legend(handles=[teal_patch, rose_patch], fontsize=7, loc='upper left', framealpha=0.4)
    fig.tight_layout()
    return _savefig(fig)


def chart_velocity():
    """Horizontal bar: traditional team vs AI pair workstreams."""
    fig, ax = plt.subplots(figsize=(5.5, 3.2), facecolor='none')
    fig.patch.set_alpha(0)
    ax.set_facecolor('none')
    streams = ['Deployment\nconfig', 'Design system', 'Copilot\nintegration', 'Simulation\nengine',
               'Personal\nfrontend', 'Advisor\nfrontend', 'JWT auth', 'Backend\n+ schema']
    trad = [2.5, 4, 4, 4, 15, 15, 4, 9]
    ai   = [0.5, 1, 1, 1, 4, 4, 1, 2]
    y = np.arange(len(streams))
    ax.barh(y + 0.2, trad, height=0.35, color=_c(_NAVY), alpha=0.7, label='Traditional team (days)')
    ax.barh(y - 0.2, ai,   height=0.35, color=_c(_BLUE), alpha=0.9, label='1 PO + Claude (days)')
    ax.set_yticks(y)
    ax.set_yticklabels(streams, fontsize=7.5)
    ax.set_xlabel('Estimated Calendar Days', fontsize=8)
    ax.set_title('Build Velocity Comparison', fontsize=10, fontweight='bold', color=_c(_NAVY), pad=8)
    ax.legend(fontsize=7.5, loc='lower right', framealpha=0.4)
    for sp in ['top', 'right']:
        ax.spines[sp].set_visible(False)
    ax.spines['left'].set_color('#AAAAAA')
    ax.spines['bottom'].set_color('#AAAAAA')
    ax.tick_params(labelsize=7.5)
    fig.tight_layout()
    return _savefig(fig)


def add_picture(slide, buf, left, top, width):
    slide.shapes.add_picture(buf, left, top, width=width)


# ── SLIDE BUILDERS ────────────────────────────────────────────────────────────

def slide_01_cover(prs):
    """Slide 1 — Cover."""
    s = blank_slide(prs)
    rect(s, 0, 0, W, H, fill_color=DARK_NAVY)
    rect(s, 0, 0, Inches(0.18), H, fill_color=BRAND_BLUE)
    rect(s, Inches(0.5), Inches(1.5), W - Inches(1.0), Inches(0.04), fill_color=BRAND_BLUE)

    # Big title
    txbox(s, Inches(0.6), Inches(1.6), Inches(10.0), Inches(1.0),
          text='A-RiA  ·  ARIA Personal',
          size=Pt(36), bold=True, color=WHITE)
    txbox(s, Inches(0.6), Inches(2.7), Inches(10.0), Inches(0.55),
          text='"Real Intelligence for Every Client"  ·  "Your Money Intelligence"',
          size=Pt(17), italic=True, color=RGBColor(0xA8, 0xC8, 0xE8))

    rect(s, Inches(0.6), Inches(3.55), Inches(5.0), Inches(0.04), fill_color=BRAND_BLUE)

    # Tagline block
    txbox(s, Inches(0.6), Inches(3.72), Inches(10.0), Inches(0.5),
          text='Two fintech products. One shared platform. Real intelligence for every investor.',
          size=Pt(13), color=RGBColor(0xCC, 0xDD, 0xEE), wrap=True)

    meta = [
        ('Products', 'A-RiA (Advisor Workbench) · ARIA Personal (Consumer App)'),
        ('Author',   'Sunny Hayes (sunder-vasudevan)'),
        ('Date',     'March 2026'),
        ('Stack',    'FastAPI · Supabase · React + Vite · Tailwind · Render · Vercel'),
    ]
    for i, (k, v) in enumerate(meta):
        txbox(s, Inches(0.6), Inches(4.4) + i * Inches(0.38),
              Inches(1.5), Inches(0.32),
              text=k.upper(), size=Pt(8), bold=True, color=TEAL)
        txbox(s, Inches(2.3), Inches(4.4) + i * Inches(0.38),
              Inches(6.5), Inches(0.32),
              text=v, size=Pt(8), color=RGBColor(0xCC, 0xDD, 0xEE))

    rect(s, 0, H - Inches(0.3), W, Inches(0.3), fill_color=RGBColor(0x00, 0x1A, 0x3A))
    txbox(s, Inches(0.3), H - Inches(0.28), W - Inches(0.6), Inches(0.26),
          text='Made with \u2764\ufe0f in Hyderabad  |  \u00a9 2026 Sunny Hayes. All rights reserved.',
          size=Pt(7), color=RGBColor(0x88, 0x99, 0xAA))


def slide_02_problem(prs):
    """Slide 2 — The Problem."""
    s = blank_slide(prs)
    section_label(s, '01  The Problem')
    slide_title(s, 'Wealth Advisory Is Broken for Most Indians')
    rule_line(s, Inches(1.2))

    # Left: advisor pain
    rect(s, Inches(0.5), Inches(1.4), Inches(5.9), Inches(4.8), fill_color=LIGHT_GREY)
    txbox(s, Inches(0.7), Inches(1.55), Inches(5.5), Inches(0.32),
          text='THE ADVISOR PROBLEM', size=Pt(9), bold=True, color=MID_BLUE)
    advisor_items = [
        'No actionable intelligence — data in spreadsheets',
        'Reactive, not proactive — calls clients after market moves',
        'Goal tracking is manual — no real-time probability',
        'Meeting prep = 15–30 mins of manual work per client',
        'Urgency is invisible — high-risk clients blend in',
    ]
    bullet_block(s, advisor_items, Inches(0.7), Inches(2.0), Inches(5.5),
                 size=Pt(10.5), spacing=Inches(0.44))

    # Right: investor pain
    rect(s, Inches(6.9), Inches(1.4), Inches(5.9), Inches(4.8), fill_color=RGBColor(0xE8, 0xF2, 0xFC))
    txbox(s, Inches(7.1), Inches(1.55), Inches(5.5), Inches(0.32),
          text='THE SELF-DIRECTED INVESTOR PROBLEM', size=Pt(9), bold=True, color=BRAND_BLUE)
    investor_items = [
        'Data without intelligence — no probability engine',
        'No personalised guidance — generic content only',
        'Life event blind spots — decisions in isolation',
        'The advisor gap — cannot afford or do not want one',
        'No answer to: "Will I actually reach my goal?"',
    ]
    bullet_block(s, investor_items, Inches(7.1), Inches(2.0), Inches(5.5),
                 size=Pt(10.5), spacing=Inches(0.44))

    callout(s,
            'ARIA addresses both groups with purpose-built products on a shared intelligent platform.',
            Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.52))

    footer_bar(s)


def slide_03_two_products(prs):
    """Slide 3 — The Two Products."""
    s = blank_slide(prs)
    section_label(s, '02  The Products')
    slide_title(s, 'Two Products. One Shared Platform.')
    rule_line(s, Inches(1.2))

    # A-RiA card (left, dark navy)
    rect(s, Inches(0.5), Inches(1.4), Inches(5.9), Inches(5.0), fill_color=DARK_NAVY)
    txbox(s, Inches(0.65), Inches(1.55), Inches(5.6), Inches(0.5),
          text='A-RiA', size=Pt(26), bold=True, color=WHITE)
    txbox(s, Inches(0.65), Inches(2.05), Inches(5.6), Inches(0.32),
          text='Advisor Relationship Intelligence Assistant', size=Pt(10), color=TEAL, italic=True)
    txbox(s, Inches(0.65), Inches(2.42), Inches(5.6), Inches(0.28),
          text='"Real Intelligence for Every Client"', size=Pt(9), bold=True, color=RGBColor(0xA8, 0xC8, 0xE8))
    rect(s, Inches(0.65), Inches(2.82), Inches(5.0), Inches(0.03), fill_color=TEAL)
    a_items = [
        'Client book with urgency scoring',
        'Client 360 — full profile + portfolio + goals',
        'Goal probability (Monte Carlo, inflation-adj.)',
        'What-if scenario — SIP sliders, reverse calc',
        'AI copilot (Claude API) — natural language',
        'Morning briefing — daily urgent client digest',
        'Audit logs for compliance',
    ]
    for i, item in enumerate(a_items):
        txbox(s, Inches(0.65), Inches(3.0) + i * Inches(0.34), Inches(5.5), Inches(0.3),
              text=f'\u25b8  {item}', size=Pt(9.5), color=RGBColor(0xC8, 0xD8, 0xF0))

    # ARIA Personal card (right, brand blue)
    rect(s, Inches(6.9), Inches(1.4), Inches(5.9), Inches(5.0), fill_color=BRAND_BLUE)
    txbox(s, Inches(7.05), Inches(1.55), Inches(5.6), Inches(0.5),
          text='ARIA Personal', size=Pt(26), bold=True, color=WHITE)
    txbox(s, Inches(7.05), Inches(2.05), Inches(5.6), Inches(0.32),
          text='Consumer App for Self-Directed Investors', size=Pt(10), color=WHITE, italic=True)
    txbox(s, Inches(7.05), Inches(2.42), Inches(5.6), Inches(0.28),
          text='"Your Money Intelligence"', size=Pt(9), bold=True, color=RGBColor(0xE0, 0xF0, 0xFF))
    rect(s, Inches(7.05), Inches(2.82), Inches(5.0), Inches(0.03), fill_color=WHITE)
    p_items = [
        'Register / Login with JWT auth (consumer-grade)',
        'Dashboard — portfolio, goals, life events',
        'Goals with Monte Carlo probability pills',
        'What-if v2 — SIP + inflation + return sliders',
        'Life events — home, education, retirement',
        'Ask ARIA — Claude API copilot, user-scoped',
        'Mobile-first layout',
    ]
    for i, item in enumerate(p_items):
        txbox(s, Inches(7.05), Inches(3.0) + i * Inches(0.34), Inches(5.65), Inches(0.3),
              text=f'\u25b8  {item}', size=Pt(9.5), color=WHITE)

    footer_bar(s)


def slide_04_architecture(prs):
    """Slide 4 — Architecture diagram (text-based)."""
    s = blank_slide(prs)
    section_label(s, '03  Architecture')
    slide_title(s, 'One Backend. Two Frontends. Shared Intelligence.')
    rule_line(s, Inches(1.2))

    # Frontend boxes
    rect(s, Inches(1.0), Inches(1.45), Inches(4.5), Inches(1.0), fill_color=DARK_NAVY)
    txbox(s, Inches(1.1), Inches(1.55), Inches(4.3), Inches(0.3),
          text='a-ria.vercel.app', size=Pt(10), bold=True, color=WHITE)
    txbox(s, Inches(1.1), Inches(1.85), Inches(4.3), Inches(0.45),
          text='A-RiA Advisor Frontend\nReact 18 + Vite + Tailwind', size=Pt(8.5), color=TEAL, wrap=True)

    rect(s, Inches(7.9), Inches(1.45), Inches(4.5), Inches(1.0), fill_color=BRAND_BLUE)
    txbox(s, Inches(8.0), Inches(1.55), Inches(4.3), Inches(0.3),
          text='aria-personal.vercel.app', size=Pt(10), bold=True, color=WHITE)
    txbox(s, Inches(8.0), Inches(1.85), Inches(4.3), Inches(0.45),
          text='ARIA Personal Frontend\nReact 18 + Vite + Tailwind', size=Pt(8.5), color=WHITE, wrap=True)

    # Arrows down
    for x in [Inches(3.1), Inches(10.15)]:
        txbox(s, x, Inches(2.5), Inches(0.5), Inches(0.4),
              text='\u2193', size=Pt(20), bold=True, color=MID_GREY, align=PP_ALIGN.CENTER)

    # Backend box
    rect(s, Inches(1.5), Inches(2.95), Inches(10.4), Inches(1.8), fill_color=MID_BLUE)
    txbox(s, Inches(1.7), Inches(3.05), Inches(10.0), Inches(0.3),
          text='aria-advisor.onrender.com  —  FastAPI Backend', size=Pt(11), bold=True, color=WHITE)
    rect(s, Inches(1.7), Inches(3.38), Inches(4.8), Inches(0.9),
         fill_color=RGBColor(0x00, 0x40, 0x80))
    txbox(s, Inches(1.8), Inches(3.45), Inches(4.6), Inches(0.28),
          text='/clients/* routes  (Advisor API)', size=Pt(9), bold=True, color=WHITE)
    txbox(s, Inches(1.8), Inches(3.73), Inches(4.6), Inches(0.48),
          text='urgency scoring · goal probability\nclient 360 · copilot', size=Pt(8.5), color=TEAL, wrap=True)

    rect(s, Inches(7.0), Inches(3.38), Inches(4.8), Inches(0.9),
         fill_color=RGBColor(0x12, 0x55, 0xB5))
    txbox(s, Inches(7.1), Inches(3.45), Inches(4.6), Inches(0.28),
          text='/personal/* routes  (Consumer API + JWT)', size=Pt(9), bold=True, color=WHITE)
    txbox(s, Inches(7.1), Inches(3.73), Inches(4.6), Inches(0.48),
          text='auth · goals · life events\ncopilot · portfolio', size=Pt(8.5), color=WHITE, wrap=True)

    # Shared row inside backend
    rect(s, Inches(1.7), Inches(4.33), Inches(10.0), Inches(0.32),
         fill_color=RGBColor(0x00, 0x1A, 0x4A))
    txbox(s, Inches(1.8), Inches(4.37), Inches(9.8), Inches(0.24),
          text='Shared:  simulation.py (Monte Carlo)  ·  Claude API (Copilot)  ·  JWT auth', size=Pt(8.5), color=TEAL)

    # Arrow down
    txbox(s, Inches(6.4), Inches(4.78), Inches(0.5), Inches(0.38),
          text='\u2193', size=Pt(20), bold=True, color=MID_GREY, align=PP_ALIGN.CENTER)

    # Database box
    rect(s, Inches(3.0), Inches(5.2), Inches(7.3), Inches(0.9), fill_color=DARK_NAVY)
    txbox(s, Inches(3.15), Inches(5.28), Inches(7.0), Inches(0.28),
          text='Supabase PostgreSQL  (pooler :6543)', size=Pt(10), bold=True, color=WHITE)
    txbox(s, Inches(3.15), Inches(5.56), Inches(7.0), Inches(0.42),
          text='clients · portfolios · holdings · goals · life_events · PersonalUser · PersonalCopilotLog', size=Pt(8), color=TEAL, wrap=True)

    footer_bar(s)


def slide_05_advisor_features(prs):
    """Slide 5 — A-RiA Advisor Features."""
    s = blank_slide(prs)
    section_label(s, '04  A-RiA Advisor')
    slide_title(s, 'A-RiA: Know Before They Call')
    rule_line(s, Inches(1.2))

    features = [
        ('Client Book + Urgency Scoring',
         'Clients ranked by algorithmic urgency — portfolio drift, goal probability, missed SIP, '
         'life events, interaction recency. The RM sees immediately who needs attention today.'),
        ('Client 360',
         'Full-page view: profile, portfolio, goals with probability pills, life events, '
         'urgency flags — everything needed for a meeting in one screen.'),
        ('Goal Probability — Monte Carlo',
         '1,000-simulation engine. Target inflated for real purchasing power. '
         'Probability pills: green \u226570%, amber 40–70%, rose <40%.'),
        ('What-If Scenario Engine',
         'Mode 1: "Will I achieve it?" — SIP sliders + inflation rate + real vs nominal corpus. '
         'Mode 2: "What SIP do I need?" — binary-search to 80% probability.'),
        ('AI Copilot (Claude API)',
         'Natural language interface. Context-aware: reads client data, drafts meeting notes, '
         'surfaces anomalies, answers questions in plain language.'),
        ('Morning Briefing + Audit Logs',
         'Daily digest of urgent clients and pending actions. Full compliance trail on all advisor actions.'),
    ]

    colors = [DARK_NAVY, MID_BLUE, BRAND_BLUE, TEAL, MID_BLUE, DARK_NAVY]

    for i, (title, desc) in enumerate(features):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.3)
        y = Inches(1.45) + row * Inches(1.65)
        rect(s, x, y, Inches(5.9), Inches(1.45),
             fill_color=LIGHT_GREY if col == 0 else RGBColor(0xE8, 0xF2, 0xFC))
        rect(s, x, y, Inches(0.18), Inches(1.45), fill_color=colors[i])
        txbox(s, x + Inches(0.28), y + Inches(0.1), Inches(5.45), Inches(0.3),
              text=title, size=Pt(10), bold=True, color=colors[i])
        txbox(s, x + Inches(0.28), y + Inches(0.44), Inches(5.45), Inches(0.88),
              text=desc, size=Pt(9.5), color=DARK_TEXT, wrap=True)

    footer_bar(s)


def slide_06_personal_features(prs):
    """Slide 6 — ARIA Personal Features."""
    s = blank_slide(prs)
    section_label(s, '05  ARIA Personal')
    slide_title(s, 'ARIA Personal: Your Money Intelligence')
    rule_line(s, Inches(1.2))

    features = [
        ('JWT Authentication',
         'Email + password register/login. python-jose + passlib, 7-day tokens, '
         'localStorage persistence. Production-grade consumer auth from day one.'),
        ('Goals + Monte Carlo Probability',
         'Set goal name, target ₹, date, monthly SIP. On save: Monte Carlo runs, probability pill appears. '
         'Real vs nominal corpus displayed. What-if v2 sliders for instant scenario replay.'),
        ('What-If v2 Sliders',
         'SIP, inflation rate, and expected return sliders with debounced auto-run. '
         'See in real time: +₹5,000/month SIP raises probability from 45% to 72%.'),
        ('Life Events',
         'Log planned milestones: home purchase, education, marriage, retirement, custom events. '
         'Linked to goals for full context in AI copilot.'),
        ('Ask ARIA — AI Copilot',
         'Claude API, user-scoped context. "Your retirement goal has 52% probability. '
         'To reach 80%, increase SIP by ₹8,400/month or extend timeline by 3 years."'),
        ('Mobile-First Design',
         '375px viewport first. Thumb-friendly navigation, cards, and forms. '
         'Safari-compatible date inputs. India\'s investors are mobile-first.'),
    ]

    colors = [BRAND_BLUE, TEAL, MID_BLUE, DARK_NAVY, BRAND_BLUE, TEAL]

    for i, (title, desc) in enumerate(features):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.3)
        y = Inches(1.45) + row * Inches(1.65)
        rect(s, x, y, Inches(5.9), Inches(1.45),
             fill_color=LIGHT_GREY if col == 0 else RGBColor(0xE8, 0xF2, 0xFC))
        rect(s, x, y, Inches(0.18), Inches(1.45), fill_color=colors[i])
        txbox(s, x + Inches(0.28), y + Inches(0.1), Inches(5.45), Inches(0.3),
              text=title, size=Pt(10), bold=True, color=colors[i])
        txbox(s, x + Inches(0.28), y + Inches(0.44), Inches(5.45), Inches(0.88),
              text=desc, size=Pt(9.5), color=DARK_TEXT, wrap=True)

    footer_bar(s)


def slide_07_simulation(prs):
    """Slide 7 — The Simulation Engine."""
    s = blank_slide(prs)
    section_label(s, '06  Simulation Engine')
    slide_title(s, 'Monte Carlo Goal Probability — The Analytical Core')
    rule_line(s, Inches(1.2))

    # Left: explanation
    txbox(s, Inches(0.5), Inches(1.4), Inches(6.5), Inches(0.3),
          text='HOW IT WORKS', size=Pt(9), bold=True, color=MID_GREY)
    steps = [
        '1.  Inflation-adjust target: real_target = amount \u00d7 (1 + rate)^years',
        '2.  For each of 1,000 paths, step month-by-month:',
        '     r \u223c N(annual_rate/12,  0.05/\u221a12)  [Gaussian return]',
        '     value = value \u00d7 (1 + r) + monthly_SIP',
        '3.  Count successes: final_value \u2265 real_target',
        '4.  Return probability_pct + median corpus (nominal & real)',
    ]
    for i, step in enumerate(steps):
        txbox(s, Inches(0.5), Inches(1.8) + i * Inches(0.42), Inches(6.2), Inches(0.38),
              text=step, size=Pt(10), color=DARK_TEXT)

    txbox(s, Inches(0.5), Inches(4.5), Inches(6.5), Inches(0.3),
          text='FIND_REQUIRED_SIP', size=Pt(9), bold=True, color=MID_GREY)
    txbox(s, Inches(0.5), Inches(4.85), Inches(6.2), Inches(0.8),
          text='Binary-search over SIP range. 30 iterations \u2192 precision within \u20b91. '
               'Rounded to nearest \u20b9100. Powers: "What SIP do I need for 80% probability?"',
          size=Pt(10), color=DARK_TEXT, wrap=True)

    # Right: charts
    try:
        mc_buf = chart_monte_carlo()
        add_picture(s, mc_buf, left=Inches(7.0), top=Inches(1.35), width=Inches(5.8))
    except Exception:
        pass

    try:
        pill_buf = chart_probability_pills()
        add_picture(s, pill_buf, left=Inches(7.3), top=Inches(4.45), width=Inches(5.5))
    except Exception:
        pass

    callout(s,
            'Why inflation-adjustment? A corpus of \u20b950 lakh in 2045 \u2260 \u20b950 lakh today. '
            'ARIA always shows a real-purchasing-power probability.',
            Inches(0.5), Inches(5.85), Inches(6.2), Inches(0.65))

    footer_bar(s)


def slide_08_design(prs):
    """Slide 8 — Design System."""
    s = blank_slide(prs)
    section_label(s, '07  Design System')
    slide_title(s, 'One Design Language. Two Products.')
    rule_line(s, Inches(1.2))

    # Colour swatches
    palette = [
        (BRAND_BLUE, '#1D6FDB', 'Brand Blue', 'Primary actions, links'),
        (DARK_NAVY,  '#002B5C', 'Dark Navy',  'Login panels, headers'),
        (TEAL,       '#007A87', 'Teal',       'On-track pills'),
        (AMBER,      '#E88C00', 'Amber',      'Needs-attention pills'),
        (ROSE,       '#E53E3E', 'Rose',       'At-risk pills'),
        (LIGHT_GREY, '#F2F4F7', 'Slate',      'Card backgrounds'),
    ]
    for i, (color, hex_val, name, usage) in enumerate(palette):
        x = Inches(0.5) + i * Inches(2.1)
        rect(s, x, Inches(1.45), Inches(1.85), Inches(0.7), fill_color=color)
        txbox(s, x, Inches(2.18), Inches(1.85), Inches(0.24),
              text=hex_val, size=Pt(8), bold=True, color=MID_GREY, align=PP_ALIGN.CENTER)
        txbox(s, x, Inches(2.42), Inches(1.85), Inches(0.26),
              text=name, size=Pt(8.5), bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
        txbox(s, x, Inches(2.68), Inches(1.85), Inches(0.28),
              text=usage, size=Pt(7.5), color=MID_GREY, align=PP_ALIGN.CENTER, wrap=True)

    # Design decisions
    rule_line(s, Inches(3.1), color=LIGHT_GREY)
    txbox(s, Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.28),
          text='KEY DESIGN DECISIONS', size=Pt(9), bold=True, color=MID_GREY)

    decisions = [
        ('Probability Pills',
         'Three-state teal/amber/rose system is the primary data visualisation. '
         'Consistent semantics: green = on track, amber = review needed, red = intervene.'),
        ('Login Layout',
         'Split-screen: left dark navy panel (tagline + stats), right slate-50 panel (form). '
         'Premium first impression, fast mobile load.'),
        ('Logo Components',
         'ARiALogo.jsx (advisor) and ARIALogo.jsx (personal) — dotless-i with round blue dot. '
         'Consistent wordmark across both products.'),
        ('Mobile-First',
         '375px viewport first. All navigation, cards, and interactive elements are thumb-friendly. '
         'India\'s investors are on mobile.'),
    ]
    for i, (title, desc) in enumerate(decisions):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.3)
        y = Inches(3.55) + row * Inches(1.35)
        rect(s, x, y, Inches(5.9), Inches(1.15),
             fill_color=LIGHT_GREY if col == 0 else RGBColor(0xE8, 0xF2, 0xFC))
        txbox(s, x + Inches(0.15), y + Inches(0.1), Inches(5.6), Inches(0.28),
              text=title, size=Pt(10), bold=True, color=BRAND_BLUE)
        txbox(s, x + Inches(0.15), y + Inches(0.42), Inches(5.6), Inches(0.65),
              text=desc, size=Pt(9.5), color=DARK_TEXT, wrap=True)

    footer_bar(s)


def slide_09_engineering(prs):
    """Slide 9 — Engineering Story: One PO + Claude."""
    s = blank_slide(prs)
    section_label(s, '08  Engineering Story')
    slide_title(s, 'One Product Owner + Claude Built Two Fintech Apps Simultaneously')
    rule_line(s, Inches(1.2))

    # The story
    txbox(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.55),
          text='A-RiA and ARIA Personal were designed, architected, coded, debugged, and deployed in a single build window by one person '
               'with Claude as a persistent full-session AI pair-programmer.',
          size=Pt(11), color=DARK_TEXT, wrap=True)

    # Challenges table
    txbox(s, Inches(0.5), Inches(2.1), Inches(12.0), Inches(0.28),
          text='KEY TECHNICAL CHALLENGES SOLVED IN THE SAME SESSION', size=Pt(9), bold=True, color=MID_GREY)

    two_col_table(s,
        headers=['Challenge', 'Fix', 'Lesson'],
        rows=[
            ('NameError: models.Goal / models.LifeEvent in clients.py',
             'Corrected model imports in clients.py',
             'FastAPI lazy eval — import errors don\'t surface until first call'),
            ('Safari <input type="date"> not firing onChange',
             'Replaced with month + year <select> elements, state lifted to parent',
             'Test on Safari. India\'s mobile users = iOS = Safari.'),
            ('FastAPI 204 delete routes returning 500',
             'Return explicit Response(status_code=204)',
             'Always be explicit on status-only responses.'),
            ('NOT NULL constraint — personal_user_id on shared tables',
             'Nullable FK + explicit query filters in both route namespaces',
             'Nullable FKs with explicit filters = clean dual-product schema.'),
        ],
        left=Inches(0.5), top=Inches(2.45),
        col_widths=[Inches(3.8), Inches(3.8), Inches(4.7)],
        row_height=Inches(0.5)
    )

    callout(s,
            'Two codebases. One context. When the Safari bug was found in the advisor app, it was fixed in the '
            'personal app in the same turn — no ticket, no second PR, no risk of inconsistency.',
            Inches(0.5), Inches(5.15), Inches(12.3), Inches(0.72))

    footer_bar(s)


def slide_10_velocity(prs):
    """Slide 10 — Velocity Metrics."""
    s = blank_slide(prs)
    section_label(s, '09  Velocity')
    slide_title(s, 'Build Velocity: One PO + Claude vs. Traditional 3-Person Team')
    rule_line(s, Inches(1.2))

    # KPI row
    kpis = [
        ('2', 'Production apps\nshipped'),
        ('3', 'Cloud services\nconfigured'),
        ('1', 'Shared simulation\nengine'),
        ('35–50', 'Features across\nboth products'),
        ('20–40x', 'Est. compression\nratio'),
    ]
    for i, (num, lbl) in enumerate(kpis):
        big_number(s, num, lbl, left=Inches(0.4) + i * Inches(2.4), top=Inches(1.5))

    rule_line(s, Inches(3.0), color=LIGHT_GREY)

    # Comparison table
    two_col_table(s,
        headers=['Workstream', 'Traditional Team (est.)', '1 PO + Claude'],
        rows=[
            ('FastAPI backend + Supabase schema', '1–2 weeks', 'Single session'),
            ('JWT auth system', '3–5 days', 'Same session'),
            ('Advisor frontend (full feature set)', '2–3 weeks', 'Same session'),
            ('Personal frontend (full feature set)', '2–3 weeks', 'Same session'),
            ('Monte Carlo simulation engine', '3–5 days', 'Same session'),
            ('TOTAL', '~8–14 weeks', '1 build window'),
        ],
        left=Inches(0.5), top=Inches(3.1),
        col_widths=[Inches(5.5), Inches(3.4), Inches(3.4)],
        row_height=Inches(0.36)
    )

    # Chart
    try:
        vel_buf = chart_velocity()
        add_picture(s, vel_buf, left=Inches(9.0), top=Inches(3.1), width=Inches(3.8))
    except Exception:
        pass

    footer_bar(s)


def slide_11_roadmap(prs):
    """Slide 11 — Roadmap."""
    s = blank_slide(prs)
    section_label(s, '10  Roadmap')
    slide_title(s, 'What\'s Next')
    rule_line(s, Inches(1.2))

    # A-RiA roadmap (left)
    rect(s, Inches(0.5), Inches(1.45), Inches(5.9), Inches(4.75), fill_color=LIGHT_GREY)
    txbox(s, Inches(0.65), Inches(1.58), Inches(5.6), Inches(0.3),
          text='A-RiA ADVISOR — PHASE 2', size=Pt(9), bold=True, color=MID_BLUE)

    aria_items = [
        ('FEAT-503', 'What-If Scenario v2', 'In design', AMBER),
        ('FEAT-301', 'Book-level copilot', 'Not started', MID_GREY),
        ('FEAT-302', 'Recommendation cards + approve/reject', 'Not started', MID_GREY),
        ('FEAT-201', 'Live NAV fetch (MFAPI.in)', 'Not started', MID_GREY),
        ('FEAT-202', 'Rebalancing proposal engine', 'Not started', MID_GREY),
        ('FEAT-101', 'Edit client data (read-only now)', 'Not started', MID_GREY),
    ]
    for i, (feat_id, name, status, color) in enumerate(aria_items):
        y = Inches(1.97) + i * Inches(0.56)
        rect(s, Inches(0.65), y, Inches(5.6), Inches(0.48),
             fill_color=WHITE if i % 2 == 0 else LIGHT_GREY)
        txbox(s, Inches(0.75), y + Inches(0.05), Inches(1.0), Inches(0.22),
              text=feat_id, size=Pt(8), bold=True, color=BRAND_BLUE)
        txbox(s, Inches(0.75), y + Inches(0.26), Inches(3.8), Inches(0.18),
              text=name, size=Pt(9), color=DARK_TEXT)
        txbox(s, Inches(4.7), y + Inches(0.1), Inches(1.3), Inches(0.28),
              text=status, size=Pt(8), italic=True, color=color, align=PP_ALIGN.RIGHT)

    # ARIA Personal roadmap (right)
    rect(s, Inches(6.9), Inches(1.45), Inches(5.9), Inches(2.45), fill_color=RGBColor(0xE8, 0xF2, 0xFC))
    txbox(s, Inches(7.05), Inches(1.58), Inches(5.6), Inches(0.3),
          text='ARIA PERSONAL — PHASE 2', size=Pt(9), bold=True, color=BRAND_BLUE)

    personal_items = [
        ('FEAT-P001', 'Portfolio add/edit UI', 'Not started'),
        ('FEAT-P002', 'Onboarding risk questionnaire', 'Not started'),
    ]
    for i, (feat_id, name, status) in enumerate(personal_items):
        y = Inches(1.97) + i * Inches(0.56)
        rect(s, Inches(7.05), y, Inches(5.6), Inches(0.48),
             fill_color=WHITE if i % 2 == 0 else LIGHT_GREY)
        txbox(s, Inches(7.15), y + Inches(0.05), Inches(1.0), Inches(0.22),
              text=feat_id, size=Pt(8), bold=True, color=BRAND_BLUE)
        txbox(s, Inches(7.15), y + Inches(0.26), Inches(3.8), Inches(0.18),
              text=name, size=Pt(9), color=DARK_TEXT)
        txbox(s, Inches(11.1), y + Inches(0.1), Inches(1.3), Inches(0.28),
              text=status, size=Pt(8), italic=True, color=MID_GREY, align=PP_ALIGN.RIGHT)

    # Platform / future
    rect(s, Inches(6.9), Inches(4.0), Inches(5.9), Inches(2.2), fill_color=DARK_NAVY)
    txbox(s, Inches(7.05), Inches(4.12), Inches(5.6), Inches(0.3),
          text='PLATFORM — FUTURE', size=Pt(9), bold=True, color=TEAL)
    future = ['Multi-currency (USD, AED — NRI segment)',
              'WhatsApp notification integration',
              'Open banking data ingestion (AA API)',
              'React Native mobile app']
    for i, item in enumerate(future):
        txbox(s, Inches(7.05), Inches(4.5) + i * Inches(0.38), Inches(5.6), Inches(0.34),
              text=f'\u25b8  {item}', size=Pt(9.5), color=RGBColor(0xC8, 0xD8, 0xF0))

    footer_bar(s)


def slide_12_cta(prs):
    """Slide 12 — Call to Action / Contact."""
    s = blank_slide(prs)
    rect(s, 0, 0, W, H, fill_color=DARK_NAVY)
    rect(s, 0, 0, Inches(0.18), H, fill_color=BRAND_BLUE)

    txbox(s, Inches(0.6), Inches(0.4), Inches(3.0), Inches(0.26),
          text='11  CALL TO ACTION', size=Pt(8), bold=True, color=TEAL)

    txbox(s, Inches(0.6), Inches(0.75), Inches(11.5), Inches(0.85),
          text='Build With Us  ·  Invest  ·  Partner', size=Pt(32), bold=True, color=WHITE)

    rect(s, Inches(0.6), Inches(1.7), Inches(7.0), Inches(0.04), fill_color=BRAND_BLUE)

    txbox(s, Inches(0.6), Inches(1.9), Inches(11.5), Inches(0.65),
          text='A-RiA and ARIA Personal are live. The Monte Carlo engine is running. The copilot is answering questions. '
               'What\'s next depends on who comes to the table.',
          size=Pt(12), color=RGBColor(0xCC, 0xDD, 0xEE), wrap=True)

    ctas = [
        ('For Advisors & Banks',
         'Pilot A-RiA with your RM team. We\'ll show you "know before they call" in your own client book.'),
        ('For Investors',
         'Try ARIA Personal. Set a goal. See your probability. Ask the copilot what it would take to improve it.'),
        ('For Founders & Engineers',
         'This is what one PO + Claude can build. Talk to us about the model, the stack, or the approach.'),
        ('For Investors',
         'Two products live. A simulation engine that\'s correct. A design system that\'s consistent. The foundation is built.'),
    ]

    colors = [TEAL, BRAND_BLUE, MID_BLUE, RGBColor(0x5B, 0x8D, 0xE0)]

    for i, (title, desc) in enumerate(ctas):
        col = i % 2
        row = i // 2
        x = Inches(0.6) + col * Inches(5.9)
        y = Inches(2.75) + row * Inches(1.5)
        rect(s, x, y, Inches(5.6), Inches(1.3), fill_color=RGBColor(0x00, 0x1A, 0x3A))
        rect(s, x, y, Inches(0.14), Inches(1.3), fill_color=colors[i])
        txbox(s, x + Inches(0.24), y + Inches(0.12), Inches(5.2), Inches(0.32),
              text=title, size=Pt(10), bold=True, color=colors[i])
        txbox(s, x + Inches(0.24), y + Inches(0.5), Inches(5.2), Inches(0.7),
              text=desc, size=Pt(9.5), color=RGBColor(0xCC, 0xDD, 0xEE), wrap=True)

    # GitHub / contact
    txbox(s, Inches(0.6), Inches(5.92), Inches(12.0), Inches(0.32),
          text='GitHub: github.com/sunder-vasudevan/aria-advisor  ·  github.com/sunder-vasudevan/aria-personal',
          size=Pt(10), color=TEAL)
    txbox(s, Inches(0.6), Inches(6.32), Inches(12.0), Inches(0.32),
          text='Live: https://a-ria.vercel.app  ·  Backend: https://aria-advisor.onrender.com',
          size=Pt(10), color=RGBColor(0xC8, 0xD8, 0xF0))

    rect(s, 0, H - Inches(0.3), W, Inches(0.3), fill_color=RGBColor(0x00, 0x0D, 0x25))
    txbox(s, Inches(0.3), H - Inches(0.28), W - Inches(0.6), Inches(0.26),
          text='Made with \u2764\ufe0f in Hyderabad  |  \u00a9 2026 Sunny Hayes (sunder-vasudevan). All rights reserved.',
          size=Pt(7), color=RGBColor(0x88, 0x99, 0xAA))


# ── Main ──────────────────────────────────────────────────────────────────────

def build_deck():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_01_cover(prs)
    slide_02_problem(prs)
    slide_03_two_products(prs)
    slide_04_architecture(prs)
    slide_05_advisor_features(prs)
    slide_06_personal_features(prs)
    slide_07_simulation(prs)
    slide_08_design(prs)
    slide_09_engineering(prs)
    slide_10_velocity(prs)
    slide_11_roadmap(prs)
    slide_12_cta(prs)

    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, 'ARIA_Executive_Deck.pptx')
    prs.save(out_path)
    print(f'\u2713  Saved: {out_path}  ({len(prs.slides)} slides)')
    return out_path


if __name__ == '__main__':
    build_deck()
